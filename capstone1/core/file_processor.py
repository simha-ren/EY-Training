"""File processing module for handling multiple file types."""
from __future__ import annotations

import os
import mimetypes
from pathlib import Path
from typing import Optional, Tuple, Dict, Any
from io import BytesIO


class FileProcessor:
    """Handle multiple file types (PDF, DOCX, CSV, XLSX, PPTX, TXT, MD)."""
    
    SUPPORTED_TYPES = {'.pdf', '.txt', '.docx', '.xlsx', '.csv', '.pptx', '.md'}
    
    @staticmethod
    def is_supported(filename: str) -> bool:
        """Check if file type is supported."""
        ext = Path(filename).suffix.lower()
        return ext in FileProcessor.SUPPORTED_TYPES
    
    @staticmethod
    def extract_text_from_pdf(file_path: str) -> str:
        """Extract text from PDF."""
        try:
            from pypdf import PdfReader
            text = []
            with open(file_path, 'rb') as f:
                pdf = PdfReader(f)
                for page in pdf.pages:
                    text.append(page.extract_text())
            return '\n'.join(text)
        except Exception as e:
            raise ValueError(f"Error processing PDF: {e}")
    
    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        """Extract text from DOCX."""
        try:
            from docx import Document
            doc = Document(file_path)
            text = []
            for para in doc.paragraphs:
                text.append(para.text)
            return '\n'.join(text)
        except Exception as e:
            raise ValueError(f"Error processing DOCX: {e}")
    
    @staticmethod
    def extract_text_from_xlsx(file_path: str) -> str:
        """Extract text from XLSX."""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path)
            text = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                text.append(f"Sheet: {sheet}")
                for row in ws.iter_rows(values_only=True):
                    text.append('\t'.join(str(v) for v in row if v is not None))
            return '\n'.join(text)
        except Exception as e:
            raise ValueError(f"Error processing XLSX: {e}")
    
    @staticmethod
    def extract_text_from_csv(file_path: str) -> str:
        """Extract text from CSV."""
        try:
            import csv
            text = []
            with open(file_path, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                for row in reader:
                    text.append(','.join(row))
            return '\n'.join(text)
        except Exception as e:
            raise ValueError(f"Error processing CSV: {e}")
    
    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        """Extract text from PPTX."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = []
            for slide_num, slide in enumerate(prs.slides, 1):
                text.append(f"Slide {slide_num}:")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return '\n'.join(text)
        except Exception as e:
            raise ValueError(f"Error processing PPTX: {e}")
    
    @staticmethod
    def extract_text_from_txt(file_path: str) -> str:
        """Extract text from TXT."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            raise ValueError(f"Error processing TXT: {e}")

    @staticmethod
    def extract_text_from_md(file_path: str) -> str:
        """Extract text from Markdown."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            # Remove markdown formatting for cleaner context
            import re
            text = re.sub(r'(!?\[.*?\]\(.*?\))', '', text)  # remove links/images
            text = re.sub(r'[#>*_`~-]+', '', text)
            text = re.sub(r'\n{2,}', '\n\n', text)
            return text.strip()
        except Exception as e:
            raise ValueError(f"Error processing Markdown: {e}")
    
    @staticmethod
    def extract_text(file_path: str) -> str:
        """Extract text from any supported file."""
        ext = Path(file_path).suffix.lower()
        
        if ext == '.pdf':
            return FileProcessor.extract_text_from_pdf(file_path)
        elif ext == '.docx':
            return FileProcessor.extract_text_from_docx(file_path)
        elif ext == '.xlsx':
            return FileProcessor.extract_text_from_xlsx(file_path)
        elif ext == '.csv':
            return FileProcessor.extract_text_from_csv(file_path)
        elif ext == '.pptx':
            return FileProcessor.extract_text_from_pptx(file_path)
        elif ext == '.txt':
            return FileProcessor.extract_text_from_txt(file_path)
        elif ext == '.md':
            return FileProcessor.extract_text_from_md(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")
    
    @staticmethod
    def get_file_metadata(file_path: str) -> Dict[str, Any]:
        """Get metadata about the file."""
        path = Path(file_path)
        return {
            "filename": path.name,
            "size_bytes": path.stat().st_size,
            "size_mb": round(path.stat().st_size / (1024 * 1024), 2),
            "extension": path.suffix.lower(),
            "created_at": path.stat().st_ctime,
            "modified_at": path.stat().st_mtime,
            "mime_type": mimetypes.guess_type(file_path)[0] or "application/octet-stream"
        }
    
    @staticmethod
    def chunk_text(text: str, chunk_size: int = 2000, overlap: int = 200) -> list:
        """Split text into overlapping chunks."""
        chunks = []
        for i in range(0, len(text), chunk_size - overlap):
            chunk = text[i:i + chunk_size]
            if chunk.strip():
                chunks.append(chunk)
        return chunks
